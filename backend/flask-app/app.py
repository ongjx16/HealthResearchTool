import os
import requests
import io
from io import BytesIO
import openai
from flask import Flask, redirect, render_template, request, url_for, jsonify, send_file
import fitz
from flask_cors import CORS, cross_origin
# from pyhtml2pdf import converter
from pptx import Presentation
from pptx.util import Inches, Pt
import json
from urllib.request import urlopen


app = Flask(__name__)
cors = CORS(app)
app.config['CORS_HEADERS']= 'Content-Type'
openai.api_key = os.getenv('OPENAI_API_KEY')




    


@app.route("/summary" )
def summary():
    input = request.args.get("input")
    model = "text-davinci-003"
    max_context_length = 4096
    input_length = len(input)
    num_chunks = (input_length // max_context_length) + 1
    chunk_size = input_length // num_chunks
    chunks = [input[i:i+chunk_size] for i in range(0, input_length, chunk_size)]
    summary = ""
    for chunk in chunks:
        response = summarise(chunk)
        summary += response
    data ={"summary": summary}
    return jsonify(data)

def summarise(input):
    # augmented_prompt = f"summarise this into multiple slides that includes header, keypoints and image_description for a presentation in json format: {input}"
    augmented_prompt = f"summarise this {input}"
    response = openai.Completion.create(
        model="text-davinci-003",
        prompt=augmented_prompt,
        temperature=.2,
        max_tokens=100,
    )["choices"][0]["text"]
    return response



@app.route("/", methods=("GET", "POST"))
@cross_origin()
def file():
    if request.method == "POST":
        file = request.files["file"]
        if not file:
            return jsonify({"error": "No file uploaded"}), 400
        try:
            pymupdf_text = splitfile(file, 1)
            model = "text-davinci-003"
            max_context_length = 4096
            input_length = len(pymupdf_text)
            num_chunks = (input_length // max_context_length) + 1
            chunk_size = input_length // num_chunks
            chunks = [pymupdf_text[i:i+chunk_size] for i in range(0, input_length, chunk_size)]
            summary = ""
            for chunk in chunks:
                response = summarise(chunk)
                summary += response
            data ={"summary": summary}
            return render_template("summary.html", summary= summary)
        except Exception as e:
            return jsonify({"error": f"Error processing file: {e}"}), 400
    return render_template("index.html")





def splitfile(file, page):
    with fitz.open(stream=file.stream.read(), filetype="pdf") as doc:
                pymupdf_text = ""
                page1 = ""
                for page in range(0, doc.page_count):
                    current_page = doc[page]
                    if page == 0:
                        page1 += current_page.get_text()
                    else:
                        pymupdf_text += current_page.get_text()
                if(page==0):
                    return(page1)
                else:
                    return(pymupdf_text)





@app.route("/chat", methods=("GET", "POST"))
def chat():
    if request.method == "POST":
        # Get the user's message from the request body
        message = request.json["message"]

        # Get the previous conversation history from the request body
        prev_messages = request.json.get("prev_messages", "")

        # Set up the OpenAI API request with all previous messages included in the prompt
        prompt = prev_messages + f"User: {message}\nAI:"
        response = openai.Completion.create(
            engine="text-davinci-002",
            prompt=prompt,
            max_tokens=1024,
            n=1,
            stop=None,
            temperature=0.5,
        )

        # Get the AI's response from the OpenAI API response
        ai_message = response.choices[0].text.strip()

        # Concatenate all conversation history into the next message prompt
        next_messages = prev_messages + f"User: {message}\nAI: {ai_message}\n"

        # Return the AI's response and the updated conversation history as a JSON object
        return jsonify({"message": ai_message, "prev_messages": next_messages})
    else:
        # Render the chat.html template for GET requests
        return render_template("chat.html")


@app.route("/generateMedia", methods=("GET", "POST"))
def index():
    input_text = None
    if request.method == "POST":
        input_text = request.form["input"]
        return redirect(url_for("generatePpt", input=input_text))
        # summarize(input_text)
    else: 
        return render_template("summarise.html")

@app.route("/<input>" )
def generatePpt(input):
    augmented_prompt = f"summarise the following text into multiple slides that includes HEADER, KEYPOINTS and IMAGE_DESCRIPTION keys to describe the image for a presentation in a json format: {input}"
    response = openai.Completion.create(
        model="text-davinci-003",
        prompt=augmented_prompt,
        temperature=.5,
        max_tokens=1000,
    )["choices"][0]["text"]
    data = '[' + response.split('[', 1)[-1]
    print(data)
    output = createPresentation(data)
    # return output
    return send_file(output, attachment_filename="presentation.pptx", as_attachment=True, cache_timeout=None)
    # return render_template("results.html", response = data )

def createPresentation(response):
    prs = Presentation()
    try:
        data = json.loads(response)
    except:
        return "error"
    for s in data:
        # sInfo = json.load(s)
        title_slide_layout = prs.slide_layouts[3]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.placeholders[0]
        title.text = s.get('HEADER')

        subtitle = slide.placeholders[1]
        print("length of placeholders")
        tf = subtitle.text_frame
        tf.margin_bottom = Inches(0.08)

        toRemove = slide.placeholders[2]
        toRemove.text =""

        for i in s.get('KEYPOINTS'):
            p = tf.add_paragraph()
            p.text = i
            font = p.font
            font.size = Pt(18)
        
        pic = slide.shapes
        # photos = pu.photos(type_='random', count=1, featured=True, query=s.get('keyword'))
        # [photo] = photos.entries
        # url = photo.link_download
        responseI = openai.Image.create(
  prompt=s.get('IMAGE_DESCRIPTION'),
  n=1,
  size="1024x1024"
)
        print("image response")
        print(responseI)
        image_url = responseI['data'][0]['url']
        image_data = BytesIO(urlopen(image_url).read())
        picAdded = pic.add_picture(image_data,Inches(5.5), Inches(1.75),width=Inches(4), height=Inches(5))
        # print(photo.id, photo.link_download)
        
        # subtitle.text = s.get('keypoints')[0]
    

    
    outputPPT = BytesIO()
    prs.save(outputPPT)
    return outputPPT




if __name__ =="__main__":
    app.run()






