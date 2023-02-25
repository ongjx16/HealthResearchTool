


from pptx import Presentation
import os

import openai
from flask import Flask, redirect, render_template, request, url_for, jsonify
import json

app = Flask(__name__)
openai.api_key = os.getenv("OPENAI_API_KEY")

@app.route("/", methods=("GET", "POST"))
def index():
    input_text = None
    if request.method == "POST":
        input_text = request.form["input"]
        return redirect(url_for("summary", input=input_text))
        # summarize(input_text)
    else: 
        return render_template("index.html")

@app.route("/<input>" )
def summary(input):
    augmented_prompt = f"summarise this into multiple slides that includes header, keypoints and image_description for a presentation in json format: {input}"
    response = openai.Completion.create(
        model="text-davinci-003",
        prompt=augmented_prompt,
        temperature=.5,
        max_tokens=1000,
    )["choices"][0]["text"]
    createPresentation(response)
    return render_template("results.html", response = response )

def createPresentation(response):
    prs = Presentation()
    data = json.loads(response)
    for s in data:
        print("s")
        print(s)
        title_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "s.header"
        subtitle.text = "s.keyPoints"
    

    prs.save('test.pptx')

